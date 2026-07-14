"""Extract text from uploaded documents and chunk it for retrieval.

Supported formats (dispatched by extension, with content-type as a fallback hint):

  * ``.txt`` / ``.md``        — decoded as UTF-8 (lenient).
  * ``.pdf``                  — ``pypdf`` (lazy import); scanned pages OCR-fallback (``ocr`` extra).
  * ``.docx``                 — ``python-docx`` (lazy import).
  * ``.pptx``                 — ``python-pptx`` (lazy import).
  * ``.xlsx``                 — ``openpyxl`` (lazy import); typed rows, one IC-13 table item per sheet.
  * ``.las``                  — ``lasio`` (lazy import); well-log curves/units/depth.

The heavy parsers are **lazy-imported inside the dispatch** so the package imports cleanly without them and a
minimal install only pays for the formats it actually uses. Reported extras: ``documents`` (``pypdf``,
``python-docx``, ``python-pptx``, ``openpyxl``, ``lasio``) and ``ocr`` (``pytesseract``, ``pdf2image``).

Chunking is a sliding window over an approximate token estimate (``~chars/4``) with overlap, which keeps related
sentences together for embedding without a tokenizer dependency. ``chunk_text`` returns plain strings;
``parse_and_chunk`` does extract-then-chunk in one call, returning flat text with no location metadata.

:func:`parse_and_chunk_located` is the location-aware sibling: it returns :class:`LocatedChunk` records that carry
page/row/col plus a content-addressed ``artifact_ref`` + ``selector`` back to the original structured data (a PDF
page image, an XLSX row range, an LAS curve set), so a retrieval hit can point at verifiable provenance instead of
only a text snippet. ``extract_typed_tables`` builds the canonical IC-13 ``KnowledgeItem`` (typed-table schema)
side of that for XLSX workbooks; see :mod:`mixle_mlops.documents.structured_artifacts`.
"""
from __future__ import annotations

import datetime as _dt
import io
from dataclasses import dataclass
from pathlib import PurePosixPath
from typing import Any

from ..multimodal.store import BlobStore, get_blob_store
from .structured_artifacts import build_typed_table_item, store_artifact

# Rough chars-per-token used to map a token budget to a character window without a tokenizer dependency.
CHARS_PER_TOKEN = 4

# A page whose extracted text is shorter than this is treated as (effectively) a scanned image and OCR'd instead.
OCR_MIN_TEXT_CHARS = 16

XLSX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


class DocumentParseError(Exception):
    """Raised when a document can't be parsed (unknown/unsupported format or a missing optional parser)."""


@dataclass
class LocatedChunk:
    """One retrievable chunk plus enough provenance to point back at canonical structure, not just text.

    ``page`` (PDF), ``row``/``col`` (XLSX) locate the chunk within its source; ``artifact_ref`` is the
    content-addressed handle (see :mod:`structured_artifacts`) to the original bytes/structured artifact, and
    ``selector`` narrows within it (e.g. ``{"page": 1}`` or ``{"sheet": "Assay", "row_start": 0, "row_end": 0,
    "columns": [...]}``). A consumer (e.g. an ``M1b``-style hydrator) uses ``artifact_ref`` + ``selector`` to
    reconstruct canonical structure instead of parsing it back out of ``text``.
    """

    text: str
    page: int | None
    row: int | None
    col: int | None
    source_format: str
    artifact_ref: str | None
    selector: dict[str, Any] | None


def _ext(filename: str) -> str:
    return PurePosixPath(filename or "").suffix.lower().lstrip(".")


def _format_of(filename: str, content_type: str | None) -> str:
    """Resolve a logical format from the filename extension, falling back to the mime type."""
    ext = _ext(filename)
    if ext in {"txt", "text", "log"}:
        return "txt"
    if ext in {"md", "markdown"}:
        return "md"
    if ext == "pdf":
        return "pdf"
    if ext in {"docx"}:
        return "docx"
    if ext in {"pptx"}:
        return "pptx"
    if ext in {"xlsx", "xlsm"}:
        return "xlsx"
    if ext == "las":
        return "las"
    ct = (content_type or "").split(";", 1)[0].strip().lower()
    mime_map = {
        "text/plain": "txt",
        "text/markdown": "md",
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
        XLSX_MEDIA_TYPE: "xlsx",
        "application/vnd.ms-excel": "xlsx",
        "application/x-las": "las",
    }
    if ct in mime_map:
        return mime_map[ct]
    if ct.startswith("text/"):
        return "txt"
    return ext or "unknown"


def _extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")


def _pdf_page_texts(data: bytes) -> list[str]:
    """Per-page text via ``pypdf`` -- ``""`` for a page whose extraction fails or is blank."""
    try:
        import pypdf  # lazy: documents extra
    except ImportError as exc:  # pragma: no cover - exercised only without the extra
        raise DocumentParseError(
            "PDF support needs the 'documents' extra (pip install pypdf)."
        ) from exc
    reader = pypdf.PdfReader(io.BytesIO(data))
    texts = []
    for page in reader.pages:
        try:
            texts.append(page.extract_text() or "")
        except Exception:
            texts.append("")
    return texts


def _extract_pdf_ocr(data: bytes, *, page: int | None = None) -> str:
    """OCR-render a PDF (or one 0-based ``page`` of it) via ``pdf2image`` + ``pytesseract`` (the ``ocr`` extra).

    Renders every page and joins the results when ``page`` is ``None`` (whole-document fallback); renders just
    that page when given (the per-page fallback used by :func:`_extract_pdf` / :func:`parse_and_chunk_located`).
    """
    try:
        from pdf2image import convert_from_bytes  # lazy: ocr extra
        import pytesseract  # lazy: ocr extra
    except ImportError as exc:
        raise DocumentParseError(
            "Scanned-PDF OCR needs the 'ocr' extra (pip install pytesseract pdf2image; also requires the "
            "poppler + tesseract system binaries)."
        ) from exc
    kwargs = {} if page is None else {"first_page": page + 1, "last_page": page + 1}
    images = convert_from_bytes(data, **kwargs)
    return "\n\n".join((pytesseract.image_to_string(im) or "") for im in images)


def _pdf_page_texts_with_ocr_fallback(data: bytes) -> list[str]:
    """Per-page text, OCR-ing (best-effort) any page whose text layer is under :data:`OCR_MIN_TEXT_CHARS`."""
    texts = _pdf_page_texts(data)
    out = []
    for i, text in enumerate(texts):
        if len(text.strip()) < OCR_MIN_TEXT_CHARS:
            try:
                text = _extract_pdf_ocr(data, page=i) or text
            except DocumentParseError:
                pass  # 'ocr' extra (or its system binaries) unavailable -- keep whatever text layer there was
        out.append(text)
    return out


def _extract_pdf(data: bytes) -> str:
    parts = _pdf_page_texts_with_ocr_fallback(data)
    return "\n\n".join(p for p in parts if p)


def _extract_docx(data: bytes) -> str:
    try:
        import docx  # python-docx, lazy: documents extra
    except ImportError as exc:  # pragma: no cover
        raise DocumentParseError(
            "DOCX support needs the 'documents' extra (pip install python-docx)."
        ) from exc
    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs if p.text]
    for table in document.tables:                          # include table cell text
        for row in table.rows:
            cells = [c.text for c in row.cells if c.text]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx, lazy: documents extra
    except ImportError as exc:  # pragma: no cover
        raise DocumentParseError(
            "PPTX support needs the 'documents' extra (pip install python-pptx)."
        ) from exc
    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text:
                        parts.append(text)
    return "\n".join(parts)


# --- XLSX: header annotation parsing + type inference -----------------------------------------------------

def _parse_header(raw: Any) -> dict[str, Any]:
    """Parse one header cell into ``{name, unit, is_pk}``.

    ``*Name`` marks a declared primary-key column; ``Name [unit]`` records a unit alongside the bare name.
    """
    text = "" if raw is None else str(raw).strip()
    is_pk = text.startswith("*")
    if is_pk:
        text = text[1:].strip()
    name, unit = text, None
    if text.endswith("]") and "[" in text:
        head, _, tail = text.rpartition("[")
        candidate_unit = tail[:-1].strip()
        if candidate_unit:
            name, unit = head.strip(), candidate_unit
    return {"name": name or "col", "unit": unit, "is_pk": is_pk}


def _infer_type(values: list[Any]) -> str:
    non_null = [v for v in values if v is not None]
    if not non_null:
        return "string"
    if all(isinstance(v, bool) for v in non_null):
        return "boolean"
    if all(isinstance(v, int) and not isinstance(v, bool) for v in non_null):
        return "integer"
    if all(isinstance(v, (int, float)) and not isinstance(v, bool) for v in non_null):
        return "float"
    if all(isinstance(v, (_dt.datetime, _dt.date)) for v in non_null):
        return "datetime"
    return "string"


def _coerce_value(value: Any, ctype: str) -> Any:
    if value is None:
        return None
    if ctype == "datetime":
        if isinstance(value, (_dt.datetime, _dt.date)):
            return value.isoformat()
        return str(value)
    if ctype == "boolean":
        return bool(value)
    if ctype == "integer":
        return int(value)
    if ctype == "float":
        return float(value)
    return value if isinstance(value, str) else str(value)


def _infer_columns_and_rows(
    col_specs: list[dict[str, Any]], raw_rows: list[list[Any]]
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    columns = []
    for i, spec in enumerate(col_specs):
        values = [r[i] if i < len(r) else None for r in raw_rows]
        ctype = _infer_type(values)
        nullable = (not spec["is_pk"]) and any(v is None for v in values)
        columns.append({"name": spec["name"], "type": ctype, "unit": spec["unit"], "nullable": nullable})
    rows = []
    for r in raw_rows:
        row = {}
        for i, spec in enumerate(col_specs):
            raw_v = r[i] if i < len(r) else None
            row[spec["name"]] = _coerce_value(raw_v, columns[i]["type"])
        rows.append(row)
    return columns, rows


def _xlsx_sheets(data: bytes) -> list[dict[str, Any]]:
    """Parse every worksheet into ``{name, columns, rows, primary_key}`` (row 1 = header)."""
    try:
        import openpyxl  # lazy: documents extra
    except ImportError as exc:
        raise DocumentParseError(
            "XLSX support needs the 'documents' extra (pip install openpyxl)."
        ) from exc
    workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheets = []
    for ws in workbook.worksheets:
        rows_iter = ws.iter_rows(values_only=False)
        header_cells = next(rows_iter, None)
        if header_cells is None:
            sheets.append({"name": ws.title, "columns": [], "rows": [], "primary_key": ["_row_id"]})
            continue
        col_specs = [_parse_header(c.value) for c in header_cells]
        raw_rows = [[c.value for c in row[: len(col_specs)]] for row in rows_iter]

        primary_key = [c["name"] for c in col_specs if c["is_pk"]]
        if not primary_key:                                # no declared key -- synthesize a sequential one
            col_specs = [{"name": "_row_id", "unit": None, "is_pk": True}, *col_specs]
            raw_rows = [[i + 1, *r] for i, r in enumerate(raw_rows)]
            primary_key = ["_row_id"]

        columns, rows = _infer_columns_and_rows(col_specs, raw_rows)
        sheets.append({"name": ws.title, "columns": columns, "rows": rows, "primary_key": primary_key})
    return sheets


def _row_text(columns: list[dict[str, Any]], row: dict[str, Any]) -> str:
    """``"col=val, col2=val2"`` rendering of a row -- an embedding/display surface only, never load-bearing."""
    return ", ".join(f"{c['name']}={row.get(c['name'])}" for c in columns)


def _extract_xlsx(data: bytes) -> str:
    sheets = _xlsx_sheets(data)
    parts = []
    for sheet in sheets:
        parts.append(f"## {sheet['name']}")
        parts.extend(_row_text(sheet["columns"], row) for row in sheet["rows"])
    return "\n".join(parts)


def extract_typed_tables(data: bytes, *, filename: str, store: BlobStore) -> list[dict[str, Any]]:
    """One IC-13 ``KnowledgeItem`` dict (``mixle://schema/typed-table/1``) per sheet of an XLSX workbook.

    The raw workbook bytes are persisted once as an immutable artifact; every returned item's ``artifact_ref``
    points at that same artifact and its ``metadata``/``relations`` carry the sheet name through for provenance.
    """
    sheets = _xlsx_sheets(data)
    artifact = store_artifact(
        data, filename=filename or "workbook.xlsx", media_type=XLSX_MEDIA_TYPE, store=store
    )
    items = []
    for sheet in sheets:
        item_id = f"typed-table:{artifact.sha256[:16]}:{sheet['name']}"
        text_surface = "\n".join(_row_text(sheet["columns"], row) for row in sheet["rows"])
        items.append(
            build_typed_table_item(
                item_id=item_id,
                sheet_name=sheet["name"],
                primary_key=sheet["primary_key"],
                columns=sheet["columns"],
                rows=sheet["rows"],
                workbook_ref=artifact.ref,
                workbook_sha256=artifact.sha256,
                workbook_media_type=XLSX_MEDIA_TYPE,
                source_filename=filename,
                text_surface=text_surface,
            )
        )
    return items


# --- LAS: well-log curves -----------------------------------------------------------------------------------

def _las_curves(data: bytes) -> dict[str, Any]:
    """Parse an LAS well log into ``{well, curves: [{mnemonic, unit, descr, data}], depth_start, depth_stop}``."""
    try:
        import lasio  # lazy: documents extra
    except ImportError as exc:
        raise DocumentParseError(
            "LAS support needs the 'documents' extra (pip install lasio)."
        ) from exc
    las = lasio.read(io.StringIO(data.decode("utf-8", errors="replace")))
    curves = [
        {
            "mnemonic": c.mnemonic,
            "unit": c.unit or None,
            "descr": c.descr or None,
            "data": [None if v != v else float(v) for v in c.data],  # v != v <=> NaN, without a numpy import
        }
        for c in las.curves
    ]
    well = {item.mnemonic: item.value for item in las.well}
    depth_values = [v for v in (curves[0]["data"] if curves else []) if v is not None]
    return {
        "well": well,
        "curves": curves,
        "depth_start": min(depth_values) if depth_values else None,
        "depth_stop": max(depth_values) if depth_values else None,
    }


def _extract_las(data: bytes) -> str:
    parsed = _las_curves(data)
    lines = []
    well_name = parsed["well"].get("WELL")
    if well_name:
        lines.append(f"Well: {well_name}")
    lines.append(
        "Curves: " + ", ".join(f"{c['mnemonic']} ({c['unit'] or 'unitless'})" for c in parsed["curves"])
    )
    if parsed["depth_start"] is not None:
        lines.append(f"Depth range: {parsed['depth_start']} - {parsed['depth_stop']}")
    return "\n".join(lines)


def extract_text(data: bytes, *, filename: str = "", content_type: str | None = None) -> str:
    """Extract plain text from ``data`` for a supported format (dispatched by extension/mime).

    PDF pages whose text layer is effectively empty are OCR-fallback'd (best-effort; silently skipped if the
    ``ocr`` extra/system binaries aren't installed).
    """
    fmt = _format_of(filename, content_type)
    if fmt in {"txt", "md"}:
        return _extract_txt(data)
    if fmt == "pdf":
        return _extract_pdf(data)
    if fmt == "docx":
        return _extract_docx(data)
    if fmt == "pptx":
        return _extract_pptx(data)
    if fmt == "xlsx":
        return _extract_xlsx(data)
    if fmt == "las":
        return _extract_las(data)
    raise DocumentParseError(
        f"unsupported document format {fmt!r} (filename={filename!r}, content_type={content_type!r})"
    )


def chunk_text(
    text: str,
    *,
    chunk_tokens: int = 256,
    overlap_tokens: int = 32,
    chunk_chars: int | None = None,
    overlap_chars: int | None = None,
) -> list[str]:
    """Sliding-window chunking with overlap.

    Sizes are given in approximate tokens (mapped to characters via :data:`CHARS_PER_TOKEN`); pass
    ``chunk_chars`` / ``overlap_chars`` to specify a character window directly. Windows are cut at the nearest
    whitespace before the boundary when possible, so chunks don't split mid-word.
    """
    text = (text or "").strip()
    if not text:
        return []
    size = chunk_chars if chunk_chars is not None else max(1, chunk_tokens * CHARS_PER_TOKEN)
    over = overlap_chars if overlap_chars is not None else max(0, overlap_tokens * CHARS_PER_TOKEN)
    over = min(over, size - 1)                              # overlap must be strictly less than the window
    chunks: list[str] = []
    n = len(text)
    start = 0
    while start < n:
        end = min(start + size, n)
        if end < n:                                        # try to break on whitespace for cleaner chunks
            window = text[start:end]
            cut = window.rfind(" ")
            if cut > size // 2:                            # only honour the break if it isn't too early
                end = start + cut
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        start = max(end - over, start + 1)
    return chunks


def parse_and_chunk(
    data: bytes,
    *,
    filename: str = "",
    content_type: str | None = None,
    chunk_tokens: int = 256,
    overlap_tokens: int = 32,
) -> tuple[str, list[str]]:
    """Extract text then chunk it. Returns ``(full_text, chunks)``."""
    text = extract_text(data, filename=filename, content_type=content_type)
    chunks = chunk_text(text, chunk_tokens=chunk_tokens, overlap_tokens=overlap_tokens)
    return text, chunks


# --- location-aware chunking --------------------------------------------------------------------------------

def _located_pdf(data: bytes, *, filename: str, tokens: int, overlap: int) -> list[LocatedChunk]:
    page_texts = _pdf_page_texts_with_ocr_fallback(data)
    artifact = store_artifact(
        data, filename=filename or "document.pdf", media_type="application/pdf", store=get_blob_store()
    )
    chunks: list[LocatedChunk] = []
    for i, text in enumerate(page_texts):
        for c in chunk_text(text, chunk_tokens=tokens, overlap_tokens=overlap):
            chunks.append(
                LocatedChunk(
                    text=c, page=i, row=None, col=None, source_format="pdf",
                    artifact_ref=artifact.ref, selector={"page": i},
                )
            )
    return chunks


def _located_xlsx(data: bytes, *, filename: str, tokens: int, overlap: int) -> list[LocatedChunk]:
    items = extract_typed_tables(data, filename=filename, store=get_blob_store())
    sheets = _xlsx_sheets(data)
    chunks: list[LocatedChunk] = []
    for item, sheet in zip(items, sheets):
        col_names = [c["name"] for c in sheet["columns"]]
        for row_index, row in enumerate(sheet["rows"]):
            chunks.append(
                LocatedChunk(
                    text=_row_text(sheet["columns"], row),
                    page=None, row=row_index, col=None, source_format="xlsx",
                    artifact_ref=item.get("artifact_ref"),
                    selector={
                        "sheet": sheet["name"], "row_start": row_index, "row_end": row_index,
                        "columns": col_names,
                    },
                )
            )
    return chunks


def _located_las(data: bytes, *, filename: str, tokens: int, overlap: int) -> list[LocatedChunk]:
    parsed = _las_curves(data)
    artifact = store_artifact(
        data, filename=filename or "log.las", media_type="application/x-las", store=get_blob_store()
    )
    curve_names = [c["mnemonic"] for c in parsed["curves"]]
    text = _extract_las(data)
    selector = {"curves": curve_names, "depth_start": parsed["depth_start"], "depth_stop": parsed["depth_stop"]}
    return [
        LocatedChunk(
            text=c, page=None, row=None, col=None, source_format="las",
            artifact_ref=artifact.ref, selector=selector,
        )
        for c in chunk_text(text, chunk_tokens=tokens, overlap_tokens=overlap)
    ]


def parse_and_chunk_located(
    data: bytes,
    *,
    filename: str = "",
    content_type: str | None = None,
    tokens: int = 400,
    overlap: int = 60,
) -> list[LocatedChunk]:
    """Extract + chunk ``data`` into :class:`LocatedChunk` records carrying page/row/col + artifact provenance.

    PDF pages get an OCR fallback (see :func:`extract_text`); XLSX sheets get one chunk per row plus a canonical
    IC-13 typed-table item per sheet (see :func:`extract_typed_tables`); LAS curves get a structured artifact.
    Any other supported format falls back to flat, location-less chunking (``page=row=col=None``).
    """
    fmt = _format_of(filename, content_type)
    if fmt == "pdf":
        return _located_pdf(data, filename=filename, tokens=tokens, overlap=overlap)
    if fmt == "xlsx":
        return _located_xlsx(data, filename=filename, tokens=tokens, overlap=overlap)
    if fmt == "las":
        return _located_las(data, filename=filename, tokens=tokens, overlap=overlap)
    text = extract_text(data, filename=filename, content_type=content_type)
    return [
        LocatedChunk(
            text=c, page=None, row=None, col=None, source_format=fmt, artifact_ref=None, selector=None,
        )
        for c in chunk_text(text, chunk_tokens=tokens, overlap_tokens=overlap)
    ]
